# -*- coding: utf-8 -*-  
import os
#import imp
import sys
if sys.version_info < (3, 0):
    reload(sys)
    sys.setdefaultencoding('utf8')
import xlrd
from xlutils.copy import copy
import re
#import xlutils
#from xlutils.copy import copy as xlutils_copy
from websocket_server import WebsocketServer
from pptx import Presentation
from docx import Document

class OpenXml:
    def __init__(self, oxml, type, list):
        self.openXml = oxml
        self.list = list
        self.type = type

    def save(self, path):
        self.openXml.save(path)

class ExcelData:
    def __init__(self, table, i, j):
        self.x = i
        self.y = j
        self.table = table

    def write(self, text):
        self.table.write(self.x, self.y, unicode(text))

class PPTData:
    def __init__(self, run):
        self.run = run

    def write(self, text):
        self.run.text = unicode(text)

class WordData:
    def __init__(self, paragraph):
        self.paragraph = paragraph

    def write(self, text):
        #print('bofore:' + self.paragraph.text)
        self.paragraph.text = unicode(text)
        #print('after:' + self.paragraph.text)

class Chain:
    class Array:
        def __init__(self):
            self.array = []

        def write(self, text):
            for meta in self.array:
                meta.write(text)

        def append(self, meta):
            self.array.append(meta)

    def __init__(self):
        self.jitsu = {}

    def addIfNew(self, text, meta):
        array = self.jitsu.get(text)
        if array == None:
            array = Chain.Array()
            array.append(meta)
            self.jitsu[text] = array
            return (False, array)

        array.append(meta)
        return (True, array) 


def isString(var):
    return (isinstance(var, unicode) or isinstance(var, str))

def isAllCharacters(var):
    return len(re.findall(r'[\d\w\s\-_\|\.\&\*\(\)\!@#$。！（）——“”？]', var)) == len(var)

def openDocument(path):
    t = splitFileType(path)
    if t == 'xls' or t == 'xlsx':
        return openXLS(path)
    elif t == 'pptx':
        return openPPTX(path)
    elif t == 'docx':
        return openDOCX(path)
    else:
        print('NOT SUPPORTED')

def openDOCX(path):
    docx = Document(path)
    paragraphs = docx.paragraphs
    compound = []
    chain = Chain()

    def add_text(pointer, compound, chain):
        text = pointer.text
        if isString(text):
            v = text.strip()
            if v != '' and not isAllCharacters(v):
                meta = WordData(pointer)
                added, array = chain.addIfNew(text, meta)
                if not added:
                    compound.append((text, array))

    for p in paragraphs:
        add_text(p, compound, chain)

    for table in docx.tables:
        for row in table.rows:
            for cell in row.cells:
                add_text(cell, compound, chain)

    return OpenXml(docx, 'docx', compound)

def openPPTX(path):
    prs = Presentation(path)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    compound = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if isString(var):
                        v = run.text.strip()
                        if v != '' and not isAllCharacters(v):
                            print(run.text)
                            #compound.append((run.text, run))
                            compound.append((run.text, PPTData(run)))

    return OpenXml(prs, 'pptx', compound)

def openXLS(path):
    book = xlrd.open_workbook(path, on_demand=True, formatting_info=True)
    wb = copy(book)
    #table = data.sheet_by_index(0)
    tables = book.sheets()
    count = len(tables)
    compound = []
    chain = Chain()


#    while i < count:
    for table in tables:
        wb_table = wb.get_sheet(table.name)
        nRows = table.nrows
        nCols = table.ncols
        for i in range(nRows):
            for j in range(nCols):
                v = table.cell(i, j).value                
                if isString(v):
                    v2 = v.strip()
                    if v2 != '' and not isAllCharacters(v2):
                        element = table.cell(i, j)
                        #compound.append((element.value, element))
                        meta = ExcelData(wb_table, i, j)
                        added, array = chain.addIfNew(element.value, meta)
                        if not added:
                            compound.append((element.value, array))
                        print(element.value)

    return OpenXml(wb, 'xlsx', compound)

def translateByChromeExtension(src):
    data = [0, src]

    def sendMessage(server, client):
        server.send_message(client, src[index])
        index += 1

    def new_client(client, server):
        print("New client connected and was given id %d" % client['id'])
        data = new_client.data
        if data[0] < len(data[1]):
            #server.send_message(client, data[1][data[0]])
            server.send_message(client, data[1][data[0]][0])
            data[0] += 1
        else:
            print('finished!')
        #server.send_message_to_all("Hey all, a new client has joined us")

    new_client.data = data

    # Called for every client disconnecting
    def client_left(client, server):
        print("Client(%d) disconnected" % client['id'])


    # Called when a client sends a message
    def message_received(client, server, message):
        #if len(message) > 200:
	    #    message = message[:200]+'..'
        #print("<== Client(%d) said: %s" % (client['id'], message))
        data = message_received.data
        
        #print('==> ' + data[1][data[0]])
        print('==>:' + data[1][data[0] - 1][0])
        print("<==:" + message)        

        #data[1][data[0] -1][1].text = message
        data[1][data[0] -1][1].write(message)
        
        if data[0] < len(data[1]):
            #server.send_message(client, data[1][data[0]])
            server.send_message(client, data[1][data[0]][0])
            data[0] += 1
        else:
            print('finished!')

    message_received.data = data

    # Called for every client connecting (after handshake)
    PORT = 9001
    server = WebsocketServer(PORT, host='127.0.0.1')#, loglevel=logging.INFO)
    server.set_fn_new_client(new_client)
    server.set_fn_client_left(client_left)
    server.set_fn_message_received(message_received)
    server.run_forever()

def splitFileType(f):
    t = f[f.rfind('.') + 1 :]
    return t

if __name__ == '__main__':
    path = sys.path[0] + '/' + sys.argv[1]
    openxml = openDocument(path)
    #source = readXLS(path)
    #if len(source) <= 0:
    #    print('data source out of range.')
    
    translateByChromeExtension(openxml.list)
    openxml.save(sys.argv[1] + '.dest.' + openxml.type)
